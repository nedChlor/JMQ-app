#!/usr/bin/env python3
"""
JMQ Service Manual — BUILD ALL
Master build script. Generates all project data from source files.

Usage:
  python src/scripts/build_all.py

Outputs:
  app_assets/              ← ready for Flutter
  src/output/              ← developer artifacts (JSON, CSV, docs)
"""

import json, csv, re, sqlite3, shutil, sys, subprocess
from pathlib import Path
from collections import defaultdict

PROJECT = Path(__file__).resolve().parent.parent.parent

# Paths
SRC_DOCS   = PROJECT / 'src' / 'data_raw' / 'documents'
SRC_XLSX   = PROJECT / 'src' / 'data_raw' / 'dtc_source.xlsx'
OUT_DEV    = PROJECT / 'src' / 'output'
OUT_APP    = PROJECT / 'app_assets'
OUT_DB     = OUT_APP / 'db'
OUT_PDF    = OUT_APP / 'pdf'

for p in [OUT_DEV, OUT_DB, OUT_PDF]:
    p.mkdir(parents=True, exist_ok=True)

print(f'JMQ Service Manual — Build All')
print(f'  Source docs:  {SRC_DOCS}')
print(f'  Source xlsx:  {SRC_XLSX}')
print(f'  Output dev:   {OUT_DEV}')
print(f'  Output app:   {OUT_APP}')
print()

# =========================================================================
# 1. CATEGORIES — predefined hierarchy
# =========================================================================
CATEGORIES = {
    '01_Engine': {
        'name_ru': 'Двигатель', 'name_en': 'Engine', 'icon': 'engine', 'order': 1,
        'sub': {
            'HFC4GB2.4D':  {'name_ru': 'HFC4GB2.4D (1.5T)', 'name_en': 'HFC4GB2.4D (1.5T)', 'order': 1},
            'HFC4GB2.4E':  {'name_ru': 'HFC4GB2.4E (1.5T Euro6)', 'name_en': 'HFC4GB2.4E (1.5T Euro6)', 'order': 2},
        }
    },
    '02_Transmission': {
        'name_ru': 'Трансмиссия', 'name_en': 'Transmission', 'icon': 'transmission', 'order': 2,
        'sub': {
            'CVT_Punch': {'name_ru': 'CVT Punch (вариатор)', 'name_en': 'CVT Punch (CVT)', 'order': 1},
            'MF622D35':  {'name_ru': 'MF622D35 (мех. КПП)', 'name_en': 'MF622D35 (Manual)', 'order': 2},
        }
    },
    '03_Chassis': {
        'name_ru': 'Шасси', 'name_en': 'Chassis', 'icon': 'chassis', 'order': 3, 'sub': {},
    },
    '04_Body': {
        'name_ru': 'Кузов', 'name_en': 'Body', 'icon': 'body', 'order': 4,
        'sub': {
            '01_Structure':  {'name_ru': 'Структура кузова', 'name_en': 'Body Structure', 'order': 1},
            '02_Electrical': {'name_ru': 'Электрооборудование', 'name_en': 'Electrical Equipment', 'order': 2},
        }
    },
    '05_Electrical_Systems': {
        'name_ru': 'Электрические системы', 'name_en': 'Electrical Systems', 'icon': 'electrical', 'order': 5,
        'sub': {
            '01_Circuit_Diagrams': {'name_ru': 'Принципиальные схемы', 'name_en': 'Circuit Diagrams', 'order': 1},
            '02_AVM':              {'name_ru': 'AVM (камера 360°)', 'name_en': 'AVM (360° Camera)', 'order': 2},
            '03_TPMS':             {'name_ru': 'TPMS (давление в шинах)', 'name_en': 'TPMS (Tire Pressure)', 'order': 3},
        }
    },
    '06_Calibration': {
        'name_ru': 'Калибровка', 'name_en': 'Calibration', 'icon': 'calibration', 'order': 6, 'sub': {},
    },
    '99_Firmware_Updates': {
        'name_ru': 'Обновления ПО', 'name_en': 'Firmware Updates', 'icon': 'firmware', 'order': 99,
        'sub': {
            'MP5': {'name_ru': 'MP5 (мультимедиа)', 'name_en': 'MP5 (Multimedia)', 'order': 1},
        }
    },
}

def build_categories(src_dir):
    cats = []
    doc_dir_map = {}  # cat_id -> absolute dir path
    cat_id = 0

    for dir_name in sorted(CATEGORIES.keys()):
        cat = CATEGORIES[dir_name]
        parent_path = src_dir / dir_name
        if not parent_path.exists():
            continue

        cat_id += 1
        parent_id = cat_id
        doc_dir_map[parent_id] = parent_path

        cats.append({
            'id': parent_id, 'name_ru': cat['name_ru'], 'name_en': cat['name_en'],
            'icon': cat['icon'], 'parent_id': None, 'sort_order': cat['order'],
        })

        for sub_name in sorted(cat['sub'].keys()):
            sub_info = cat['sub'][sub_name]
            sub_path = parent_path / sub_name
            if not sub_path.exists():
                continue

            cat_id += 1
            doc_dir_map[cat_id] = sub_path
            cats.append({
                'id': cat_id, 'name_ru': sub_info['name_ru'], 'name_en': sub_info['name_en'],
                'icon': cat['icon'], 'parent_id': parent_id, 'sort_order': sub_info['order'],
            })

    return cats, doc_dir_map


# =========================================================================
# 2. DOCUMENT TEXT EXTRACTION
# =========================================================================
def extract_pdf(path):
    try:
        import pdfplumber
        text = []
        with pdfplumber.open(str(path)) as pdf:
            for p in pdf.pages:
                t = p.extract_text()
                if t: text.append(t)
        return '\n\n'.join(text)
    except Exception as e:
        return f'[Extraction error: {e}]'

def extract_docx(path):
    try:
        import docx
        doc = docx.Document(str(path))
        paras = [p.text for p in doc.paragraphs if p.text.strip()]
        for t in doc.tables:
            for r in t.rows:
                rt = ' | '.join(c.text for c in r.cells if c.text.strip())
                if rt: paras.append(rt)
        return '\n'.join(paras)
    except Exception as e:
        return f'[Extraction error: {e}]'

def extract_pptx(path):
    try:
        import pptx
        prs = pptx.Presentation(str(path))
        slides = []
        for i, slide in enumerate(prs.slides, 1):
            parts = [f'[Slide {i}]']
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = para.text.strip()
                        if t: parts.append(t)
                if shape.has_table:
                    for r in shape.table.rows:
                        rt = ' | '.join(c.text for c in r.cells if c.text.strip())
                        if rt: parts.append(rt)
            if len(parts) > 1:
                slides.append('\n'.join(parts))
        return '\n\n'.join(slides)
    except Exception as e:
        return f'[Extraction error: {e}]'

EXTRACTORS = {'.pdf': extract_pdf, '.docx': extract_docx, '.pptx': extract_pptx}

KNOWN_TERMS = {
    'Dvigatel': 'Двигатель', 'Transmissiya': 'Трансмиссия', 'Shassi': 'Шасси',
    'Kuzov': 'Кузов', 'Elektrika': 'Электрика', 'Kalibrovka': 'Калибровка',
    'Obnovlenie': 'Обновление', 'Rukovodstvo': 'Руководство', 'Obsluzhivanie': 'Обслуживание',
    'Mekhanika': 'Механика', 'Elektrokontrol': 'Электронный контроль',
    'Proverka': 'Проверка', 'Problem': 'Неисправности', 'Training': 'Обучение',
    'Service': 'Сервис', 'Manual': 'Руководство', 'Sistema': 'Система',
    'Obzor': 'Обзор', 'Krugovoy': 'Круговой', 'Davlenie': 'Давление', 'Shin': 'Шин',
    'Principialnye': 'Принципиальные', 'Shemy': 'Схемы', 'Elektrooborudovanie': 'Электрооборудование',
    'Struktura': 'Структура', 'Adaptaciya': 'Адаптация', 'Sopostavlenie': 'Сопоставление',
    'PO': 'ПО', 'MP5': 'MP5',
}

def make_title(filename):
    stem = filename.stem
    parts = stem.split('_')
    if len(parts) >= 2 and re.match(r'\d+\.\d+', parts[0]):
        parts = parts[1:]
    result = []
    for p in parts:
        result.append(KNOWN_TERMS.get(p, p))
    return ' '.join(result)

def build_documents(src_dir, cat_dir_map):
    out = []
    doc_id = 0
    for cat_id, dir_path in sorted(cat_dir_map.items(), key=lambda x: x[0]):
        for filepath in sorted(dir_path.iterdir()):
            ext = filepath.suffix.lower()
            if ext not in EXTRACTORS:
                continue
            doc_id += 1
            sys.stdout.write(f'  [{doc_id:2d}] {filepath.relative_to(src_dir.parent.parent)} ... ')
            sys.stdout.flush()
            text = EXTRACTORS[ext](filepath)
            print(f'{len(text):,} chars')
            out.append({
                'id': doc_id, 'category_id': cat_id,
                'title_ru': make_title(filepath),
                'title_en': filepath.stem,
                'file_type': ext[1:],
                'original_filename': filepath.name,
                'relative_path': str(filepath.relative_to(src_dir)),
                'content_text': text,
                'text_length': len(text),
                'file_size': filepath.stat().st_size,
            })
    return out


# =========================================================================
# 3. DTC EXTRACTION
# =========================================================================
import openpyxl
DTC_RE = re.compile(r'^(?:0x[0-9A-Fa-f]+|[PBCU][0-9A-Fa-f]{4,7})$')
CJK_RE = re.compile(r'[\u4e00-\u9fff\u3400-\u4dbf\uf900-\ufaff]+')

def split_cjk(text):
    if not text:
        return text, ''
    cjk_blocks = CJK_RE.findall(text)
    clean = CJK_RE.sub('', text).strip()
    clean = re.sub(r' {2,}', ' ', clean)
    cn = ' '.join(cjk_blocks).strip()
    return clean, cn

def split_ecu_name(raw):
    if not raw:
        return None, None
    raw = raw.strip()
    base_overrides = {'ESC_KM': 'ESP', 'Tbox': 'TBOX', 'TBOX': 'TBOX', 'ABS': 'ABS'}
    if raw in base_overrides:
        return base_overrides[raw], None
    raw_upper = raw.upper()
    if raw_upper in base_overrides:
        return base_overrides[raw_upper], None
    m = re.match(r'^([A-Za-z0-9_]+)([\u4e00-\u9fff\u3000-\u303f\uff00-\uffef]+)$', raw)
    if m:
        return m.group(1).upper(), m.group(2)
    m = re.match(r'^([A-Za-z0-9_]+)\s+(.+)$', raw)
    if m:
        base = m.group(1).upper()
        variant = m.group(2).strip()
        v_alpha = re.sub(r'[\u4e00-\u9fff]+', '', variant).strip()
        if v_alpha:
            return base, v_alpha
        return base, variant
    cleaned = re.sub(r'[\u4e00-\u9fff\u3000-\u303f\uff00-\uffef]+', '', raw).strip()
    return cleaned.upper() or raw.upper(), None

def parse_sheet(name):
    n = name.strip().lower()
    for hp in ('4de', 'isf', 'b4.'):
        if n.startswith(hp):
            return None, None
    direct = {'J7 PLUS': 'J7PLUS'}
    if name.strip() in direct:
        return direct[name.strip()], None
    prefix_map = {
        'j7-': 'J7', 'j7plus': 'J7PLUS', 'js4-': 'JS4', 'js7-': 'JS7',
        'js8-': 'JS8', 'js6': 'JS6', 'k7-': 'K7', 'sunray-': 'Sunray', 'es2-': 'eS2',
    }
    for p, m in prefix_map.items():
        if n.startswith(p):
            rest = name[len(p):].strip()
            return (m, rest or None)
    singles = {'S3': 'S3', 'S5': 'S5', 'J4': 'J4', 'T8': 'T8',
               'E10X': 'E10X', 'E40X': 'E40X', 'IC5': 'IC5'}
    if name.strip() in singles:
        return singles[name.strip()], None
    return None, None

def extract_sheet(ws, sheet, model, ecu_hint):
    results = []
    max_row, max_col = ws.max_row, ws.max_column

    has_ecu_col = False
    if ecu_hint is None:
        for r in range(1, 3):
            a = ws.cell(r, 1).value
            if a and ('ECU' in str(a).upper() or '\u63a7\u5236\u5668' in str(a)):
                has_ecu_col = True
                break
        if not has_ecu_col:
            for r in range(2, min(6, max_row + 1)):
                a = ws.cell(r, 1).value
                if a and DTC_RE.match(str(a).strip()):
                    break
                if a and not str(a).strip().isdigit() and len(str(a).strip()) > 1:
                    has_ecu_col = True
                    break

    code_col = meaning_col = None
    for r in range(1, min(6, max_row + 1)):
        for c in range(1, min(13, max_col + 1)):
            v = str(ws.cell(r, c).value or '').upper()
            if 'DTC DISPLAY' in v or 'DTC \u663e\u793a\u7801' in v: code_col = c
            if 'DTC MEANING' in v or 'DTC\u542b\u4e49' in v: meaning_col = c
    if code_col is None:
        for c in range(1, 6):
            for r in range(2, min(20, max_row + 1)):
                v = ws.cell(r, c).value
                if v and DTC_RE.match(str(v).strip()): code_col = c; break
            if code_col: break
    if code_col is None: return results
    if meaning_col is None: meaning_col = code_col + 1

    data_start = 5 if str(ws.cell(1, 1).value or '').startswith('Diagnostic Specification') else 2

    merged = {}
    if has_ecu_col:
        for mr in ws.merged_cells.ranges:
            if mr.min_col == 1 and mr.max_col == 1:
                v = ws.cell(mr.min_row, 1).value
                if v and not str(v).strip().isdigit():
                    for rr in range(mr.min_row, mr.max_row + 1):
                        merged[rr] = str(v).strip()

    current_ecu_raw = ecu_hint
    for r in range(data_start, max_row + 1):
        if has_ecu_col:
            if r in merged:
                current_ecu_raw = merged[r]
            else:
                a = ws.cell(r, 1).value
                if a and not str(a).strip().isdigit() and not DTC_RE.match(str(a).strip()):
                    current_ecu_raw = str(a).strip()
        if current_ecu_raw is None:
            continue

        code = ws.cell(r, code_col).value
        if not code or not DTC_RE.match(str(code).strip()):
            continue
        code = str(code).strip().upper()

        meaning_raw = ws.cell(r, meaning_col).value
        if not meaning_raw:
            meaning_raw = ws.cell(r, meaning_col + 1).value
        if not meaning_raw and meaning_col > 1:
            meaning_raw = ws.cell(r, meaning_col - 1).value
        meaning_raw = ' '.join(str(meaning_raw).split()) if meaning_raw else ''
        meaning_en, _ = split_cjk(meaning_raw)
        if not meaning_en:
            # try one more column to the right
            if meaning_col + 2 <= max_col:
                meaning_raw2 = ws.cell(r, meaning_col + 2).value
                meaning_raw2 = ' '.join(str(meaning_raw2).split()) if meaning_raw2 else ''
                meaning_en, _ = split_cjk(meaning_raw2)
        if not meaning_en:
            continue

        ecu, ecu_variant = split_ecu_name(current_ecu_raw)
        results.append({
            'model': model, 'ecu': ecu, 'ecu_variant': ecu_variant or '',
            'code': code, 'dtc_type': 'H' if code[:2] == '0X' else code[0],
            'meaning_en': meaning_en,
        })
    return results

def build_dtc(xlsx_path):
    wb = openpyxl.load_workbook(str(xlsx_path), data_only=True)
    all_dtcs = []
    for sheet in wb.sheetnames:
        if sheet == 'Catalog': continue
        model, ecu_hint = parse_sheet(sheet)
        if model is None: continue
        rows = extract_sheet(wb[sheet], sheet, model, ecu_hint)
        all_dtcs.extend(rows)

    # Group by (model, code, meaning_en), merge ecus and variants
    groups = {}
    for d in all_dtcs:
        key = (d['model'], d['code'], d['meaning_en'])
        pair = (d['ecu'], d.get('ecu_variant', ''))
        if key in groups:
            groups[key]['_pairs'].add(pair)
        else:
            d['_pairs'] = {pair}
            groups[key] = d

    result = []
    for d in groups.values():
        pairs = d.pop('_pairs')
        ecus = sorted(set(p[0] for p in pairs))
        variants = sorted(set(p[1] for p in pairs if p[1]))
        d['ecu'] = ', '.join(ecus)
        d['ecu_variant'] = ','.join(variants)
        result.append(d)
    return result


# =========================================================================
# 4. BUILD COMBINED SQLite DATABASE
# =========================================================================
def build_db(documents, dtc_codes, categories):
    db_name = 'jmq_service_manual_v5.db'
    db_path = OUT_DB / db_name
    
    if db_path.exists():
        db_path.unlink()

    conn = sqlite3.connect(str(db_path))
    conn.isolation_level = None
    conn.execute('PRAGMA journal_mode=WAL')
    conn.execute('PRAGMA synchronous=NORMAL')

    conn.execute('''CREATE TABLE vehicles (
        id INTEGER PRIMARY KEY, code TEXT UNIQUE NOT NULL, name_ru TEXT DEFAULT ''
    )''')
    conn.execute('''CREATE TABLE categories (
        id INTEGER PRIMARY KEY, name_ru TEXT NOT NULL, name_en TEXT NOT NULL,
        icon TEXT, parent_id INTEGER, sort_order INTEGER DEFAULT 0
    )''')
    conn.execute('''CREATE TABLE documents (
        id INTEGER PRIMARY KEY, category_id INTEGER NOT NULL,
        title_ru TEXT, title_en TEXT, file_type TEXT NOT NULL,
        original_filename TEXT, relative_path TEXT, content_text TEXT NOT NULL,
        text_length INTEGER DEFAULT 0, file_size INTEGER DEFAULT 0,
        FOREIGN KEY (category_id) REFERENCES categories(id)
    )''')
    conn.execute('''CREATE TABLE dtc_codes (
        id INTEGER PRIMARY KEY, vehicle_model TEXT NOT NULL,
        ecu TEXT NOT NULL, ecu_variant TEXT DEFAULT '',
        code TEXT NOT NULL, dtc_type TEXT NOT NULL, meaning_en TEXT NOT NULL
    )''')

    models = sorted(set(d['model'] for d in dtc_codes))
    for m in models:
        conn.execute('INSERT OR IGNORE INTO vehicles (code) VALUES (?)', (m,))

    for c in categories:
        conn.execute('INSERT INTO categories VALUES (?,?,?,?,?,?)',
                      (c['id'], c['name_ru'], c['name_en'], c.get('icon'),
                       c.get('parent_id'), c.get('sort_order', 0)))

    data = [(d['id'], d['category_id'], d['title_ru'], d['title_en'],
             d['file_type'], d['original_filename'], d['relative_path'],
             d['content_text'], d['text_length'], d['file_size'])
            for d in documents]
    conn.executemany('INSERT INTO documents VALUES (?,?,?,?,?,?,?,?,?,?)', data)

    data = [(d['model'], d['ecu'], d.get('ecu_variant', ''), d['code'], d['dtc_type'], d['meaning_en'])
            for d in dtc_codes]
    conn.executemany(
        'INSERT INTO dtc_codes (vehicle_model, ecu, ecu_variant, code, dtc_type, meaning_en) VALUES (?,?,?,?,?,?)',
        data)

    conn.execute('CREATE INDEX idx_doc_category ON documents(category_id)')
    conn.execute('CREATE INDEX idx_doc_type ON documents(file_type)')
    conn.execute('CREATE INDEX idx_dtc_model ON dtc_codes(vehicle_model)')
    conn.execute('CREATE INDEX idx_dtc_ecu ON dtc_codes(ecu)')
    conn.execute('CREATE INDEX idx_dtc_code ON dtc_codes(code)')
    conn.execute('CREATE INDEX idx_dtc_type ON dtc_codes(dtc_type)')
    conn.execute('CREATE INDEX idx_dtc_model_ecu ON dtc_codes(vehicle_model, ecu)')
    conn.execute('CREATE INDEX idx_dtc_variant ON dtc_codes(ecu_variant)')

    conn.execute('''CREATE VIRTUAL TABLE documents_fts USING fts5(
        title_ru, title_en, content_text,
        content='documents', content_rowid='id', tokenize='unicode61'
    )''')
    conn.execute('''CREATE VIRTUAL TABLE dtc_fts USING fts5(
        code, meaning_en, ecu, ecu_variant, vehicle_model,
        content='dtc_codes', content_rowid='id', tokenize='unicode61'
    )''')
    conn.execute("INSERT INTO documents_fts (rowid, title_ru, title_en, content_text) SELECT id, title_ru, title_en, content_text FROM documents")
    conn.execute("INSERT INTO dtc_fts (rowid, code, meaning_en, ecu, ecu_variant, vehicle_model) SELECT id, code, meaning_en, ecu, ecu_variant, vehicle_model FROM dtc_codes")
    conn.execute("INSERT INTO documents_fts(documents_fts) VALUES('optimize')")
    conn.execute("INSERT INTO dtc_fts(dtc_fts) VALUES('optimize')")

    doc_count = conn.execute('SELECT COUNT(*) FROM documents').fetchone()[0]
    dtc_count = conn.execute('SELECT COUNT(*) FROM dtc_codes').fetchone()[0]
    cat_count  = conn.execute('SELECT COUNT(*) FROM categories').fetchone()[0]
    veh_count  = conn.execute('SELECT COUNT(*) FROM vehicles').fetchone()[0]

    conn.close()
    print(f'\n  DB: {db_path}')
    print(f'  Vehicles: {veh_count}, Categories: {cat_count}, Documents: {doc_count}, DTC: {dtc_count}')
    return db_path, db_name


# =========================================================================
# 5. COPY PDFs TO APP ASSETS
# =========================================================================
def copy_pdfs():
    count = 0
    total_mb = 0
    if OUT_PDF.exists(): shutil.rmtree(str(OUT_PDF))
    for filepath in SRC_DOCS.rglob('*'):
        if filepath.suffix.lower() == '.pdf':
            rel = filepath.relative_to(SRC_DOCS).parent
            dest = OUT_PDF / rel
            dest.mkdir(parents=True, exist_ok=True)
            shutil.copy2(str(filepath), str(dest / filepath.name))
            count += 1
            total_mb += filepath.stat().st_size / 1024 / 1024
    print(f'  PDFs copied: {count} ({total_mb:.1f} MB) to {OUT_PDF}')


# =========================================================================
# 6. GENERATE DEVELOPER OUTPUTS
# =========================================================================
def write_json(path, data):
    with open(path, 'w', encoding='utf-8') as f:
        json.dump(data, f, ensure_ascii=False, indent=2)

def write_developer_outputs(categories, documents, dtc_codes, db_path=None):
    write_json(OUT_APP / 'categories.json', categories)

    meta = []
    for d in documents:
        meta.append({k: d[k] for k in ['id', 'category_id', 'title_ru', 'title_en',
                                        'file_type', 'original_filename', 'text_length', 'file_size']})
    write_json(OUT_APP / 'documents_meta.json', meta)

    enriched = []
    for d in dtc_codes:
        entry = {
            'model': d['model'], 'ecu': d['ecu'], 'ecu_variant': d.get('ecu_variant', ''),
            'code': d['code'], 'type': d['dtc_type'], 'meaning_en': d['meaning_en'],
        }
        enriched.append(entry)

    write_json(OUT_DEV / 'dtc_all_models.json', enriched)

    dtc_json_dir = OUT_DEV / 'json'; dtc_json_dir.mkdir(exist_ok=True)
    by_model = defaultdict(list)
    for d in enriched:
        by_model[d['model']].append(d)
    for model, codes in sorted(by_model.items()):
        write_json(dtc_json_dir / f'dtc_{model}.json', codes)

    csv_dir = OUT_DEV / 'reference'; csv_dir.mkdir(exist_ok=True)
    for model, codes in sorted(by_model.items()):
        path = csv_dir / f'dtc_{model}_reference.csv'
        with open(path, 'w', encoding='utf-8-sig', newline='') as f:
            w = csv.writer(f)
            w.writerow(['ECU', 'Variant', 'Code', 'Type', 'Meaning'])
            for c in codes:
                w.writerow([c['ecu'], c.get('ecu_variant', ''), c['code'], c['type'], c['meaning_en']])

    md_path = OUT_DEV / 'DTC_Reference_Manual.md'
    with open(md_path, 'w', encoding='utf-8') as f:
        f.write('# JMQ DTC Diagnostic Codes — Reference Manual\n\n')
        f.write(f'> Total: **{len(enriched)}** codes | **{len(by_model)}** models\n\n')
        f.write('| Model | Codes |\n|------|------|\n')
        for m in sorted(by_model.keys()):
            f.write(f'| {m} | {len(by_model[m])} |\n')
        f.write('\n## ECUs\n\n')
        all_ecus = sorted(set((d['ecu'], d.get('ecu_variant', '')) for d in enriched))
        for ecu, variant in all_ecus:
            label = ecu + (f' ({variant})' if variant else '')
            f.write(f'### {label}\n\n')
            for m in sorted(by_model.keys()):
                mc = [d for d in enriched if d['model'] == m and d['ecu'] == ecu and d.get('ecu_variant', '') == variant]
                if mc:
                    f.write(f'**{m}** ({len(mc)} codes)\n\n')
                    f.write('| Code | Meaning |\n|------|--------|\n')
                    for c in mc[:5]:
                        f.write(f'| {c["code"]} | {c["meaning_en"][:80]} |\n')
                    if len(mc) > 5:
                        f.write(f'| ... | +{len(mc)-5} more |\n')
                    f.write('\n')

    print(f'\nDeveloper outputs:')
    for f in sorted(OUT_DEV.rglob('*')):
        if f.is_file():
            print(f'  {f.relative_to(OUT_DEV)}')


# =========================================================================
# MAIN
# =========================================================================
print('1. Building categories...')
categories, cat_dir_map = build_categories(SRC_DOCS)

print('\n2. Extracting document text...')
documents = build_documents(SRC_DOCS, cat_dir_map)

print('\n3. Extracting DTC codes...')
dtc_codes = []
if SRC_XLSX.exists():
    dtc_codes = build_dtc(SRC_XLSX)
    print(f'   Total DTC codes: {len(dtc_codes)}')
else:
    print(f'   WARNING: DTC source file not found at {SRC_XLSX}')
    print(f'   Place dtc_source.xlsx in src/data_raw/ and re-run build_all.py')

print('\n4. Building combined SQLite database...')
db_path, db_name = build_db(documents, dtc_codes, categories)

print('\n5. Copying PDF files to app_assets...')
copy_pdfs()

print('\n6. Writing developer & app outputs...')
write_developer_outputs(categories, documents, dtc_codes, db_path=db_path)

print('\n7. Building DTC<->document cross-reference...')
crossref_script = Path(__file__).resolve().parent / 'build_crossref.py'
subprocess.run([sys.executable, str(crossref_script)], check=True)

total_chars = sum(d['text_length'] for d in documents)
total_pdf_mb = sum(d['file_size'] for d in documents if d['file_type'] == 'pdf') / 1024 / 1024
total_src_mb = sum(d['file_size'] for d in documents) / 1024 / 1024

print(f'\n{"="*50}')
print(f'BUILD COMPLETE')
print(f'{"="*50}')
print(f'  Documents:        {len(documents)}')
print(f'  Total text:       {total_chars:,} chars')
print(f'  Source files:     {total_src_mb:.0f} MB')
print(f'  PDF for app:      {total_pdf_mb:.0f} MB')
print(f'  DTC codes:        {len(dtc_codes)}')
vehicle_models = len(set(d['model'] for d in dtc_codes))
print(f'  Vehicle models:   {vehicle_models}')
print(f'')
print(f'  app_assets/  — {sum(f.stat().st_size for f in OUT_APP.rglob("*") if f.is_file()) / 1024 / 1024:.0f} MB for Flutter')
print(f'  src/output/  — developer artifacts')
print(f'  src/scripts/ — build tools')
